#!/usr/bin/env python3
"""Genera onboarding-dinamico-beglobal.pptx — versión storytelling."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG = RGBColor(0x07, 0x10, 0x1C)
PANEL = RGBColor(0x0D, 0x1A, 0x2A)
PANEL_SOFT = RGBColor(0x10, 0x1F, 0x31)
TEXT = RGBColor(0xF3, 0xF8, 0xFA)
MUTED = RGBColor(0x97, 0xAE, 0xBD)
MUTED2 = RGBColor(0x6C, 0x81, 0x90)
CYAN = RGBColor(0x55, 0xD6, 0xE8)
TEAL = RGBColor(0x18, 0xA6, 0xB9)
AMBER = RGBColor(0xFF, 0xBD, 0x66)
GREEN = RGBColor(0x70, 0xD8, 0xAA)
LINE = RGBColor(0x24, 0x35, 0x47)

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
blank = prs.slide_layouts[6]
TOTAL = 14
_n = 0


def new_slide():
    global _n
    _n += 1
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    tf = s.shapes.add_textbox(Inches(0.6), Inches(7.02), Inches(12.1), Inches(0.35)).text_frame
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "● Be Global Pro · Onboarding Dinámico"
    r.font.size = Pt(10); r.font.color.rgb = MUTED2
    tf2 = s.shapes.add_textbox(Inches(11.7), Inches(7.02), Inches(1.0), Inches(0.35)).text_frame
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run(); r2.text = f"{_n} / {TOTAL}"
    r2.font.size = Pt(10); r2.font.color.rgb = MUTED2
    return s


def run(p, text, size, color, bold=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = "Helvetica Neue"
    return r


def center_text(slide, y, h, lines, w=Inches(11.9), x=Inches(0.72)):
    """lines: list of (text_or_runs, size, color, bold, space_after)."""
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True
    first = True
    for spec in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.CENTER
        text, size, color, bold, after = spec
        p.space_after = Pt(after)
        if isinstance(text, list):
            for t, c, b in text:
                run(p, t, size, c, b)
        else:
            run(p, text, size, color, bold)
    return tf


def kicker(slide, text, y=Inches(0.55)):
    center_text(slide, y, Inches(0.4), [(text.upper(), 13, CYAN, True, 0)])


def card(slide, x, y, w, h, border=None, fill=PANEL):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.09
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border or LINE
    sh.line.width = Pt(1.4 if border else 1)
    sh.shadow.inherit = False
    return sh


def icon_card(slide, x, y, w, h, emoji, title_txt, desc, border):
    card(slide, x, y, w, h, border=border)
    tf = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.28), w - Inches(0.4), h - Inches(0.5)).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(8)
    run(p, emoji, 36, TEXT)
    q = tf.add_paragraph(); q.alignment = PP_ALIGN.CENTER; q.space_after = Pt(6)
    run(q, title_txt, 17, TEXT, True)
    d = tf.add_paragraph(); d.alignment = PP_ALIGN.CENTER
    run(d, desc, 12.5, MUTED)


def pill_row(slide, y, pills, size=13):
    total_w = 12.0
    n = len(pills)
    gap = 0.25
    pw = (total_w - gap * (n - 1)) / n
    x = (13.333 - total_w) / 2
    for txt in pills:
        c = card(slide, Inches(x), y, Inches(pw), Inches(0.62), fill=PANEL_SOFT)
        tfp = c.text_frame
        tfp.word_wrap = True
        tfp.vertical_anchor = MSO_ANCHOR.MIDDLE
        tfp.margin_top = tfp.margin_bottom = Pt(2)
        p = tfp.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run(p, txt, size, MUTED)
        x += pw + gap


def bubble(slide, x, y, w, h, text_runs, is_user):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.22
    sh.fill.solid()
    sh.fill.fore_color.rgb = TEAL if is_user else PANEL_SOFT
    sh.line.color.rgb = TEAL if is_user else LINE
    sh.line.width = Pt(1)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(12)
    tf.margin_top = tf.margin_bottom = Pt(8)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    for t, c, b in text_runs:
        run(p, t, 13, c, b)


# ── 1 Portada ──────────────────────────────────────────────
s = new_slide()
center_text(s, Inches(1.15), Inches(1.0), [("🧭", 60, TEXT, False, 0)])
kicker(s, "Be Global Pro", y=Inches(2.35))
center_text(s, Inches(2.8), Inches(1.9), [
    ("Tu guía te espera.", 46, TEXT, True, 4),
    ("Onboarding dinámico", 46, CYAN, True, 0),
])
tf = s.shapes.add_textbox(Inches(1.67), Inches(4.9), Inches(10), Inches(1.3)).text_frame
tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run(p, "La historia de cómo un socio pasa de ", 17, MUTED)
run(p, "«no sé por dónde empezar»", 17, TEXT, True)
run(p, " a ", 17, MUTED)
run(p, "su primera misión completada", 17, TEXT, True)
run(p, " — conversando, sin tecnicismos.", 17, MUTED)

# ── 2 Conoce a María ───────────────────────────────────────
s = new_slide()
center_text(s, Inches(0.75), Inches(0.95), [("👩🏽‍💼", 54, TEXT, False, 0)])
kicker(s, "Capítulo 1", y=Inches(1.85))
center_text(s, Inches(2.3), Inches(0.9), [
    ([("Conoce a ", TEXT, True), ("María", AMBER, True)], 38, TEXT, True, 0),
])
tf = s.shapes.add_textbox(Inches(2.0), Inches(3.3), Inches(9.3), Inches(1.0)).text_frame
tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run(p, "María es socia de Be Global. Tiene un producto que ama, vio todos los cursos… ", 16, MUTED)
run(p, "y aun así está atorada.", 16, TEXT, True)
pill_row(s, Inches(4.6), ["😵‍💫 Demasiado contenido, poca dirección", "🎬 Se atora con guiones y reels"], size=13)
pill_row(s, Inches(5.45), ["🛒 Quiere su tienda sin volverse técnica", "⏰ Poco tiempo y poca confianza"], size=13)

# ── 3 El deseo ─────────────────────────────────────────────
s = new_slide()
kicker(s, "Lo que María realmente quiere", y=Inches(1.0))
center_text(s, Inches(1.55), Inches(1.7), [
    ("Pedir con sus palabras.", 40, TEXT, True, 4),
    ("Recibir un camino.", 40, GREEN, True, 0),
])
card(s, Inches(3.1), Inches(3.9), Inches(7.1), Inches(2.1), border=TEAL, fill=PANEL_SOFT)
tf = s.shapes.add_textbox(Inches(3.4), Inches(4.15), Inches(6.5), Inches(1.7)).text_frame
tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(4)
run(p, "«Ayúdame a crear mi tienda.»", 20, TEXT, True)
q = tf.add_paragraph(); q.alignment = PP_ALIGN.CENTER; q.space_after = Pt(10)
run(q, "«Dime qué video debo grabar.»", 20, TEXT, True)
d = tf.add_paragraph(); d.alignment = PP_ALIGN.CENTER
run(d, "Sin aprender MCPs, APIs ni arquitectura de agentes.", 13, MUTED2)

# ── 4 Aparece la Guía ──────────────────────────────────────
s = new_slide()
center_text(s, Inches(1.3), Inches(1.0), [("✨", 56, TEXT, False, 0)])
kicker(s, "Capítulo 2", y=Inches(2.5))
center_text(s, Inches(3.0), Inches(0.95), [
    ([("Aparece la ", TEXT, True), ("Guía Be Global Pro", CYAN, True)], 38, TEXT, True, 0),
])
tf = s.shapes.add_textbox(Inches(2.2), Inches(4.15), Inches(8.9), Inches(1.2)).text_frame
tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run(p, "Un asistente que conoce la metodología, el contenido y las plantillas de Be Global. Y que empieza con ", 16, MUTED)
run(p, "una sola pregunta.", 16, TEXT, True)

# ── 5 La pregunta mágica ───────────────────────────────────
s = new_slide()
kicker(s, "Todo comienza así")
center_text(s, Inches(1.0), Inches(0.9), [
    ([("Una pregunta. ", TEXT, True), ("Tu ruta.", CYAN, True)], 36, TEXT, True, 0),
])
bubble(s, Inches(2.3), Inches(2.15), Inches(7.2), Inches(1.15), [
    ("👋 ¡Hola! Para asignarte la ruta correcta… ", MUTED, False),
    ("¿cuál es tu área en Be Global?", TEXT, True),
    ("  🏛️ Dirección · 🤝 Equipo · 🚀 Socio", MUTED, False),
], is_user=False)
bubble(s, Inches(6.0), Inches(3.55), Inches(5.0), Inches(0.75), [
    ("Soy socia 🚀 quiero vender mi producto", TEXT, False),
], is_user=True)
bubble(s, Inches(2.3), Inches(4.55), Inches(7.2), Inches(1.15), [
    ("¡Perfecto, María! Tu ruta es ", MUTED, False),
    ("Socio/Miembro", TEXT, True),
    (". Primero: ¿cuánto tiempo tienes esta semana?", MUTED, False),
], is_user=False)
center_text(s, Inches(6.0), Inches(0.6), [
    ("Nada de cuestionarios eternos. Nada de adivinar tu perfil.", 15, MUTED, False, 0),
])

# ── 6 Tres caminos ─────────────────────────────────────────
s = new_slide()
kicker(s, "Cada quien su camino")
center_text(s, Inches(1.0), Inches(0.9), [
    ([("Tres perfiles, ", TEXT, True), ("tres experiencias", GREEN, True)], 36, TEXT, True, 0),
])
icon_card(s, Inches(0.85), Inches(2.25), Inches(3.7), Inches(3.1), "🏛️", "Corporate",
          "Define el método, entrena a la Guía y gobierna la experiencia.", CYAN)
icon_card(s, Inches(4.82), Inches(2.25), Inches(3.7), Inches(3.1), "🤝", "Equipo",
          "Prueba, documenta y acompaña. Es el puente entre el método y el socio.", AMBER)
icon_card(s, Inches(8.79), Inches(2.25), Inches(3.7), Inches(3.1), "🚀", "Miembro",
          "Vive una experiencia simple y guiada hasta su primera misión.", GREEN)
center_text(s, Inches(5.75), Inches(0.6), [
    ("Cada uno en su espacio, con sus datos separados. 🔒", 15, MUTED, False, 0),
])

# ── 7 El viaje de María ────────────────────────────────────
s = new_slide()
kicker(s, "Capítulo 3 · El viaje de María")
center_text(s, Inches(1.0), Inches(0.9), [
    ([("Cinco pasos hasta ", TEXT, True), ("su primera victoria", CYAN, True)], 34, TEXT, True, 0),
])
steps = [
    ("✅", "Precheck", "Todo listo para empezar", CYAN),
    ("📝", "Intake", "Su marca, su producto, su meta", CYAN),
    ("🔌", "Setup", "Conexiones sin dolor", CYAN),
    ("🎯", "Misión 1", "Tarea concreta con entregable", CYAN),
    ("🏆", "¡Logrado!", "Evidencia en mano", GREEN),
]
x = 0.75
cw = 2.25
gap = 0.16
for emoji, t, d, c in steps:
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + cw/2 - 0.5), Inches(2.5), Inches(1.0), Inches(1.0))
    circ.fill.solid(); circ.fill.fore_color.rgb = PANEL
    circ.line.color.rgb = c; circ.line.width = Pt(2)
    circ.shadow.inherit = False
    tfc = circ.text_frame; tfc.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tfc.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run(p, emoji, 26, TEXT)
    tf = s.shapes.add_textbox(Inches(x), Inches(3.7), Inches(cw), Inches(1.6)).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(4)
    run(p, t, 16, TEXT, True)
    q = tf.add_paragraph(); q.alignment = PP_ALIGN.CENTER
    run(q, d, 11.5, MUTED2)
    if t != "¡Logrado!":
        ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + cw - 0.06), Inches(2.97), Inches(gap + 0.12), Pt(2.5))
        ln.fill.solid(); ln.fill.fore_color.rgb = TEAL
        ln.line.fill.background(); ln.shadow.inherit = False
    x += cw + gap

# ── 8 Cómo se siente ───────────────────────────────────────
s = new_slide()
kicker(s, "Así se siente cada paso")
center_text(s, Inches(1.0), Inches(1.5), [
    ("Una pregunta a la vez.", 34, TEXT, True, 4),
    ("Máximo 3 acciones.", 34, CYAN, True, 0),
])
bubble(s, Inches(2.3), Inches(2.9), Inches(7.2), Inches(1.25), [
    ("RUTA SOCIO · MISIÓN 1 · PASO 2 DE 4\n", CYAN, True),
    ("Tu guion del reel va tomando forma 🎬 Solo falta: ", MUTED, False),
    ("¿cuál es el beneficio #1 de tu producto?", TEXT, True),
], is_user=False)
bubble(s, Inches(6.0), Inches(4.4), Inches(5.0), Inches(0.75), [
    ("Que es 100% natural y dura todo el día", TEXT, False),
], is_user=True)
bubble(s, Inches(2.3), Inches(5.4), Inches(7.2), Inches(1.15), [
    ("¡Eso es oro! ✨ Aquí está tu guion con plan de tomas. Siguiente acción: ", MUTED, False),
    ("grabar la toma 1 (30 seg)", TEXT, True),
    (". ¿La agendamos para mañana?", MUTED, False),
], is_user=False)

# ── 9 Evidencia ────────────────────────────────────────────
s = new_slide()
center_text(s, Inches(0.7), Inches(0.9), [("📸", 50, TEXT, False, 0)])
kicker(s, "La regla de oro", y=Inches(1.75))
center_text(s, Inches(2.2), Inches(1.6), [
    ("Nada se da por hecho.", 34, TEXT, True, 4),
    ("Todo se demuestra.", 34, GREEN, True, 0),
])
tf = s.shapes.add_textbox(Inches(2.0), Inches(4.0), Inches(9.3), Inches(1.1)).text_frame
tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run(p, "Cada tarea se cierra con ", 15, MUTED)
run(p, "responsable, fecha y evidencia real", 15, TEXT, True)
run(p, " — un archivo, un enlace, una captura. Si algo es delicado, ", 15, MUTED)
run(p, "una persona lo revisa primero.", 15, TEXT, True)
pill_row(s, Inches(5.35), ["👤 Responsable", "📅 Fecha", "📎 Evidencia", "✅ Criterio de aceptación"], size=13)

# ── 10 Escalera de valor ───────────────────────────────────
s = new_slide()
kicker(s, "Detrás de escena")
center_text(s, Inches(1.0), Inches(0.9), [
    ([("El valor ", TEXT, True), ("sube y baja", CYAN, True), (" la escalera", TEXT, True)], 34, TEXT, True, 0),
])
rungs = [
    ("🏛️", "Corporate define", "el método y los límites", CYAN, 1.9),
    ("🤝", "Team lo domina", "lo prueba y lo documenta", AMBER, 2.5),
    ("🧑‍🏫", "Team guía al socio", "con tareas simples", AMBER, 3.1),
    ("📈", "Los resultados regresan", "para mejorar el método", GREEN, 3.7),
]
x = 0.9
cw = 2.75
gap = 0.2
base_y = 6.5
for emoji, t, d, c, hh in rungs:
    card(s, Inches(x), Inches(base_y - hh), Inches(cw), Inches(hh), border=c)
    tf = s.shapes.add_textbox(Inches(x + 0.15), Inches(base_y - hh + 0.25), Inches(cw - 0.3), Inches(hh - 0.4)).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(6)
    run(p, emoji, 28, TEXT)
    q = tf.add_paragraph(); q.alignment = PP_ALIGN.CENTER; q.space_after = Pt(4)
    run(q, t, 15, TEXT, True)
    dd = tf.add_paragraph(); dd.alignment = PP_ALIGN.CENTER
    run(dd, d, 11.5, MUTED2)
    x += cw + gap

# ── 11 Seguridad ───────────────────────────────────────────
s = new_slide()
center_text(s, Inches(0.8), Inches(0.95), [("🔐", 52, TEXT, False, 0)])
kicker(s, "Confianza ante todo", y=Inches(1.9))
center_text(s, Inches(2.4), Inches(1.6), [
    ("Tu cuenta es tuya.", 36, TEXT, True, 4),
    ("Punto.", 36, CYAN, True, 0),
])
pill_row(s, Inches(4.5), ["🙅 Nunca pide contraseñas ni códigos", "🔑 Conexiones seguras por OAuth"], size=13.5)
pill_row(s, Inches(5.35), ["🧑‍⚖️ Lo sensible pasa por aprobación humana", "🚪 Nadie ve los datos de otro miembro"], size=13.5)

# ── 12 Números ─────────────────────────────────────────────
s = new_slide()
kicker(s, "¿Cómo sabremos que funciona?")
center_text(s, Inches(1.0), Inches(0.9), [
    ([("El éxito tiene ", TEXT, True), ("números", GREEN, True)], 36, TEXT, True, 0),
])
stats = [
    ("<30 min", "del onboarding a tu primer resultado útil"),
    ("80%", "de diagnósticos de fase correctos"),
    ("8/10", "de satisfacción mínima de los socios piloto"),
]
x = 0.85
cw = 3.7
for n, l in stats:
    card(s, Inches(x), Inches(2.3), Inches(cw), Inches(2.6))
    tf = s.shapes.add_textbox(Inches(x + 0.2), Inches(2.75), Inches(cw - 0.4), Inches(2.0)).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(10)
    run(p, n, 44, CYAN, True)
    q = tf.add_paragraph(); q.alignment = PP_ALIGN.CENTER
    run(q, l, 14, MUTED)
    x += cw + 0.27
center_text(s, Inches(5.4), Inches(0.8), [
    ([("No basta con que se vea impresionante: tiene que ", MUTED, False),
      ("demostrarlo con evidencia.", TEXT, True)], 15, MUTED, False, 0),
])

# ── 13 Final ───────────────────────────────────────────────
s = new_slide()
center_text(s, Inches(0.95), Inches(1.0), [("🎉", 56, TEXT, False, 0)])
kicker(s, "Capítulo final", y=Inches(2.1))
center_text(s, Inches(2.6), Inches(1.7), [
    ("María publicó su reel.", 34, TEXT, True, 6),
    ("Y ya sabe cuál es su siguiente paso.", 30, AMBER, True, 0),
])
tf = s.shapes.add_textbox(Inches(2.0), Inches(4.7), Inches(9.3), Inches(1.2)).text_frame
tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run(p, "Eso es el onboarding dinámico: no un manual, sino una ", 16, MUTED)
run(p, "conversación que te lleva a resultados", 16, TEXT, True)
run(p, " — a ti, a tu equipo y a toda la comunidad.", 16, MUTED)

# ── 14 CTA ─────────────────────────────────────────────────
s = new_slide()
center_text(s, Inches(0.6), Inches(0.85), [("🚀", 46, TEXT, False, 0)])
kicker(s, "¿Y ahora?", y=Inches(1.5))
center_text(s, Inches(1.95), Inches(0.9), [
    ([("El piloto arranca ", TEXT, True), ("esta semana", CYAN, True)], 34, TEXT, True, 0),
])
cta = [
    ("1️⃣", "Hoy", "Confirmamos participantes y enviamos el formulario de onboarding.", CYAN),
    ("2️⃣", "En 48 h", "Kickoff de 60 min y elegimos el canal (Telegram 💙).", AMBER),
    ("3️⃣", "En 72 h", "El primer perfil vivo y los primeros escenarios probados.", GREEN),
]
x = 0.85
for emoji, t, d, c in cta:
    icon_card(s, Inches(x), Inches(3.0), Inches(3.7), Inches(2.5), emoji, t, d, c)
    x += 3.97
center_text(s, Inches(5.85), Inches(0.8), [
    ([("La Guía te ayuda a ejecutar con claridad. ", MUTED, False),
      ("No garantiza ventas ni sustituye la revisión humana.", TEXT, True)], 13.5, MUTED, False, 0),
])

out = "/Users/rogergv/Documents/SoftvibesLab/BeGlobal/beglobal/presentacion-onboarding/onboarding-dinamico-beglobal.pptx"
prs.save(out)
print("OK:", out)
